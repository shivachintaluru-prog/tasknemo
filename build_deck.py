"""Generate a TaskNemo pitch deck as PowerPoint."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Brand colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MUTED = RGBColor(0x99, 0x99, 0xAA)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
CARD_BG = RGBColor(0x25, 0x25, 0x40)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color=CARD_BG, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = 0.0
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, items, left, top, width, font_size=20, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing
        p.level = 0
    return txBox


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
         "TaskNemo", font_size=60, color=ACCENT, bold=True)
add_text(slide, Inches(1), Inches(3.0), Inches(11), Inches(1),
         "Your tasks, automatically extracted from Teams, email, calendar & meeting transcripts.",
         font_size=24, color=LIGHT_GRAY)
add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
         "AI-powered task pipeline  \u00b7  Priority scoring  \u00b7  Obsidian dashboard",
         font_size=18, color=MUTED)

# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
         "The Problem", font_size=40, color=ACCENT, bold=True)

problems = [
    "\U0001f4e8  Tasks are buried across Teams chats, emails, calendar invites, and meeting transcripts",
    "\U0001f914  You forget commitments you made in meetings \u2014 no one writes them all down",
    "\U0001f504  You manually scan channels and inbox every day to figure out what matters",
    "\u23f3  Follow-ups slip through the cracks \u2014 you only notice when someone escalates",
    "\U0001f4ca  No single view of everything you owe and everything others owe you",
]
add_bullet_slide(slide, problems, Inches(1), Inches(1.8), Inches(11), font_size=22)

# ============================================================
# SLIDE 3: What is TaskNemo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
         "What is TaskNemo?", font_size=40, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.6), Inches(11), Inches(0.8),
         "An automated pipeline that extracts your tasks from every communication channel,\n"
         "scores them by priority, and renders a live dashboard in Obsidian.",
         font_size=22, color=LIGHT_GRAY)

# Three pillars
pillars = [
    ("Extract", "Scans Teams, email, calendar,\nand meeting transcripts\nautomatically via WorkIQ"),
    ("Prioritize", "Scores tasks by stakeholder\nweight, urgency signals,\nage, and escalation patterns"),
    ("Track", "Live Obsidian dashboard\nwith state machine lifecycle\nand smart alerts"),
]

for i, (title, desc) in enumerate(pillars):
    x = Inches(1 + i * 4)
    add_shape_bg(slide, x, Inches(3.0), Inches(3.5), Inches(3.2))
    add_text(slide, x + Inches(0.3), Inches(3.3), Inches(3), Inches(0.6),
             title, font_size=28, color=ACCENT, bold=True)
    add_text(slide, x + Inches(0.3), Inches(4.0), Inches(3), Inches(2),
             desc, font_size=18, color=LIGHT_GRAY)

# ============================================================
# SLIDE 4: How It Works
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
         "How It Works", font_size=40, color=ACCENT, bold=True)

steps = [
    ("1", "Query", "Claude Code queries WorkIQ\nfor Teams, email, calendar,\ntranscripts, sent items"),
    ("2", "Extract", "AI extracts structured tasks\nfrom raw messages \u2014 both\ninbound and outbound"),
    ("3", "Dedup", "Cross-source matching by\nsender, topic, and thread ID\nprevents duplicates"),
    ("4", "Score", "Priority scoring based on\nstakeholder weight, urgency,\nage, and escalation"),
    ("5", "Render", "Dashboard in Obsidian\nwith Focus, Open, Waiting,\nand Closed sections"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.5 + i * 2.5)
    add_shape_bg(slide, x, Inches(2.0), Inches(2.3), Inches(4.0))
    add_text(slide, x + Inches(0.2), Inches(2.2), Inches(1), Inches(0.6),
             num, font_size=36, color=ACCENT2, bold=True)
    add_text(slide, x + Inches(0.2), Inches(2.8), Inches(2), Inches(0.6),
             title, font_size=22, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.2), Inches(3.5), Inches(2), Inches(2.2),
             desc, font_size=16, color=LIGHT_GRAY)

# ============================================================
# SLIDE 5: Sources
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
         "What It Scans", font_size=40, color=ACCENT, bold=True)

sources = [
    ("\U0001f4ac Teams", "Channel messages, group chats,\n1:1 DMs \u2014 extracts asks\nand commitments"),
    ("\U0001f4e7 Email", "Inbox emails, flagged items,\n@mentions in documents,\nreply status tracking"),
    ("\U0001f4c5 Calendar", "Meeting invites with agendas,\naction items from notes,\nattendee tracking"),
    ("\U0001f399\ufe0f Transcripts", "Full meeting transcripts \u2014\nthe richest source for\nhidden commitments"),
    ("\U0001f4e4 Sent Items", "Your own replies and\ndeliverables \u2014 prevents\nfalse positive tasks"),
    ("\U0001f50d Outbound", "Messages you sent where\nrecipient hasn't replied \u2014\nauto \"Waiting on Others\""),
]

for i, (title, desc) in enumerate(sources):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.2)
    y = Inches(1.8 + row * 2.8)
    add_shape_bg(slide, x, y, Inches(3.8), Inches(2.4))
    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.3), Inches(0.6),
             title, font_size=22, color=GREEN, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.3), Inches(1.2),
             desc, font_size=16, color=LIGHT_GRAY)

# ============================================================
# SLIDE 6: Dashboard
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
         "The Dashboard", font_size=40, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(5), Inches(0.6),
         "Live in Obsidian \u2014 check off tasks right in the doc",
         font_size=18, color=MUTED)

sections = [
    ("\U0001f3af Focus Now", "Top 5 highest-priority tasks.\nScore \u2265 70, sorted by urgency.\nWhat you should work on RIGHT NOW.", ACCENT),
    ("\u23f0 Due in 48h", "Tasks with approaching deadlines.\nSorted by due date.\nDon't let these slip.", ORANGE),
    ("\U0001f4cb Open", "All active inbound tasks,\ngrouped by source channel.\nYour full backlog at a glance.", WHITE),
    ("\u23f3 Waiting on Others", "Outbound tasks where you're\nwaiting for a reply. Nudge\nalerts when idle > 3 days.", ACCENT2),
    ("\u2705 Recently Closed", "Tasks completed in the\nlast 7 days. Proof of\nyour throughput.", GREEN),
]

for i, (title, desc, color) in enumerate(sections):
    y = Inches(2.2 + i * 1.0)
    add_shape_bg(slide, Inches(1), y, Inches(11), Inches(0.9))
    add_text(slide, Inches(1.3), y + Inches(0.1), Inches(3), Inches(0.7),
             title, font_size=20, color=color, bold=True)
    add_text(slide, Inches(5), y + Inches(0.15), Inches(6.5), Inches(0.7),
             desc.replace("\n", "  \u00b7  "), font_size=15, color=LIGHT_GRAY)

# ============================================================
# SLIDE 7: Scoring
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
         "Smart Priority Scoring", font_size=40, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
         "Every task gets a 0\u2013100 score. Higher = more urgent. Here's what drives it:",
         font_size=20, color=LIGHT_GRAY)

score_items = [
    ("Stakeholder Weight", "0\u201340 pts", "Tasks from your manager or skip-level score higher than peer requests"),
    ("Urgency Signals", "0\u201330 pts", "Keywords like \"ASAP\", \"blocker\", \"EOD\" in title or description"),
    ("Age Penalty", "0\u201320 pts", "Older tasks get boosted \u2014 things shouldn't sit for weeks"),
    ("Thread Intensity", "0\u201310 pts", "Tasks mentioned across multiple messages score higher"),
    ("Escalation Pattern", "0\u201315 pts", "Repeated mentions with increasing urgency = escalation"),
    ("User Pin", "+20 pts", "Pin tasks you want to keep in Focus Now regardless of score"),
]

for i, (name, pts, desc) in enumerate(score_items):
    col = i % 2
    row = i // 2
    x = Inches(1 + col * 6)
    y = Inches(2.4 + row * 1.5)
    add_shape_bg(slide, x, y, Inches(5.5), Inches(1.3))
    add_text(slide, x + Inches(0.3), y + Inches(0.15), Inches(3.5), Inches(0.5),
             name, font_size=20, color=WHITE, bold=True)
    add_text(slide, x + Inches(4.0), y + Inches(0.15), Inches(1.2), Inches(0.5),
             pts, font_size=18, color=ACCENT, bold=True, alignment=PP_ALIGN.RIGHT)
    add_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(5), Inches(0.5),
             desc, font_size=14, color=MUTED)

# ============================================================
# SLIDE 8: State Machine
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
         "Task Lifecycle", font_size=40, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
         "Tasks move through states automatically \u2014 no manual triage needed.",
         font_size=20, color=LIGHT_GRAY)

states = [
    ("Open", "Active task,\nyou're working on it", GREEN),
    ("Waiting", "Blocked on\nsomeone's input", ORANGE),
    ("Needs\nFollow-up", "Stale > 3 days,\nmay need a nudge", RGBColor(0xEF, 0x44, 0x44)),
    ("Likely\nDone", "Completion signal\nreceived", ACCENT2),
    ("Closed", "Terminal \u2014\ntask is done", MUTED),
]

for i, (name, desc, color) in enumerate(states):
    x = Inches(0.8 + i * 2.5)
    add_shape_bg(slide, x, Inches(2.8), Inches(2.2), Inches(2.5))
    add_text(slide, x + Inches(0.2), Inches(3.0), Inches(1.8), Inches(1.0),
             name, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(4.0), Inches(1.8), Inches(1.0),
             desc, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between states
    if i < len(states) - 1:
        add_text(slide, x + Inches(2.2), Inches(3.5), Inches(0.5), Inches(0.5),
                 "\u2192", font_size=28, color=MUTED, alignment=PP_ALIGN.CENTER)

auto_rules = [
    "No activity for 3 days \u2192 Needs Follow-up",
    "Completion signal (\"done\", \"shipped\", \"thanks\") \u2192 Likely Done",
    "Likely Done + 3 days quiet \u2192 Closed",
    "Needs Follow-up + 14 days \u2192 Auto-closed (safety net)",
    "Check [x] in Obsidian \u2192 Closed immediately",
]
add_bullet_slide(slide, auto_rules, Inches(1), Inches(5.6), Inches(11),
                 font_size=16, color=MUTED, spacing=Pt(4))

# ============================================================
# SLIDE 9: User Interaction
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
         "How You Interact", font_size=40, color=ACCENT, bold=True)

interactions = [
    ("\u2705 Check off in Obsidian", "Click the checkbox next to any task \u2014 it closes on next sync.\nNo CLI needed for the most common action."),
    ("\U0001f4e5 Task Inbox", "Drop tasks into Task Inbox.md in your vault.\nThey're automatically imported on next refresh."),
    ("\U0001f4cc Pin important tasks", "python task_dashboard.py pin TASK-042\nInstant +20 score boost, stays in Focus Now."),
    ("\U0001f916 Talk to Claude", "\"Add a task to follow up with Jane on the proposal\"\nClaude creates it with the right metadata."),
    ("\u23f0 Scheduled sync", "Runs automatically every 2 hours during work hours.\nDesktop toast notifications for new tasks and alerts."),
]

for i, (title, desc) in enumerate(interactions):
    y = Inches(1.6 + i * 1.15)
    add_shape_bg(slide, Inches(1), y, Inches(11.3), Inches(1.0))
    add_text(slide, Inches(1.3), y + Inches(0.05), Inches(3.5), Inches(0.5),
             title, font_size=20, color=WHITE, bold=True)
    add_text(slide, Inches(5), y + Inches(0.1), Inches(7), Inches(0.8),
             desc, font_size=15, color=LIGHT_GRAY)

# ============================================================
# SLIDE 10: Setup
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
         "Getting Started", font_size=40, color=ACCENT, bold=True)

add_text(slide, Inches(1), Inches(1.6), Inches(11), Inches(0.6),
         "5-minute setup. No admin rights needed.", font_size=22, color=LIGHT_GRAY)

setup_steps = [
    ("1. Clone & Install", "git clone <repo> && cd tasknemo\n.\\install.ps1\n\nSets up data files, configures WorkIQ MCP,\nasks for your Obsidian vault path."),
    ("2. First Sync", "Open Claude Code in the repo folder.\nSay: \"Run a full TaskNemo sync\"\n\nClaude queries all your channels and\nbuilds your initial task list."),
    ("3. Schedule", ".\\setup_full_sync_scheduler.ps1\n\nFull sync every 2 hours (weekdays 9\u20137).\nLightweight refresh every 30 min.\nDesktop toasts for new tasks."),
]

for i, (title, desc) in enumerate(setup_steps):
    x = Inches(1 + i * 4)
    add_shape_bg(slide, x, Inches(2.8), Inches(3.5), Inches(3.8))
    add_text(slide, x + Inches(0.3), Inches(3.0), Inches(3), Inches(0.6),
             title, font_size=24, color=GREEN, bold=True)
    add_text(slide, x + Inches(0.3), Inches(3.7), Inches(3), Inches(2.8),
             desc, font_size=16, color=LIGHT_GRAY)

# ============================================================
# SLIDE 11: Prerequisites
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
         "Prerequisites", font_size=40, color=ACCENT, bold=True)

prereqs = [
    ("Required", [
        "Python 3.10+",
        "Windows 11",
        "Claude Code CLI (npm install -g @anthropic-ai/claude-code)",
        "WorkIQ MCP (configured by install script)",
    ]),
    ("Recommended", [
        "Obsidian (for viewing the dashboard \u2014 any markdown editor works too)",
    ]),
    ("Automatic", [
        "Stakeholder config \u2014 auto-populated on first sync via WorkIQ lookup",
        "Scheduled sync \u2014 optional, one-time PowerShell setup",
        "Desktop notifications \u2014 pip install winotify (installed by setup)",
    ]),
]

y_pos = Inches(1.8)
for category, items in prereqs:
    add_text(slide, Inches(1), y_pos, Inches(3), Inches(0.5),
             category, font_size=24, color=GREEN, bold=True)
    y_pos += Inches(0.5)
    for item in items:
        add_text(slide, Inches(1.5), y_pos, Inches(10), Inches(0.4),
                 "\u2022  " + item, font_size=18, color=LIGHT_GRAY)
        y_pos += Inches(0.4)
    y_pos += Inches(0.3)

# ============================================================
# SLIDE 12: Why TaskNemo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
         "Why TaskNemo?", font_size=40, color=ACCENT, bold=True)

comparisons = [
    ("vs. Manual tracking", "You don't have to remember to write down every commitment from every meeting.\nTaskNemo extracts them from transcripts automatically."),
    ("vs. Planner / To Do", "Those only track what you manually add. TaskNemo finds tasks you didn't\nknow you had \u2014 buried in chats, emails, and meeting transcripts."),
    ("vs. Copilot recaps", "Copilot summarizes conversations. TaskNemo extracts actionable tasks,\nscores priority, tracks state over time, and alerts you to stale items."),
]

for i, (title, desc) in enumerate(comparisons):
    y = Inches(1.8 + i * 1.7)
    add_shape_bg(slide, Inches(1), y, Inches(11.3), Inches(1.4))
    add_text(slide, Inches(1.5), y + Inches(0.15), Inches(3.5), Inches(0.5),
             title, font_size=22, color=ACCENT, bold=True)
    add_text(slide, Inches(1.5), y + Inches(0.65), Inches(10), Inches(0.7),
             desc, font_size=17, color=LIGHT_GRAY)

# ============================================================
# SLIDE 13: CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
         "Stop losing tasks in the noise.", font_size=48, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
         "Try TaskNemo \u2014 5 minutes to set up, runs on autopilot after that.",
         font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
         "git clone <repo>  \u00b7  .\\install.ps1  \u00b7  \"Run a full TaskNemo sync\"",
         font_size=20, color=MUTED, alignment=PP_ALIGN.CENTER)

# Save
output_path = os.path.join(os.path.dirname(__file__), "TaskNemo_Pitch.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
